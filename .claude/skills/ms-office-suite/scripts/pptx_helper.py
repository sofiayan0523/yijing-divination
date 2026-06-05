#!/usr/bin/env python3
"""PowerPoint helper — create, read, and edit .pptx files.

Usage:
    python pptx_helper.py create <output.pptx> <title> [subtitle]
    python pptx_helper.py read <input.pptx>
    python pptx_helper.py info <input.pptx>
"""
import sys
import json


def create_presentation(output_path: str, title: str, subtitle: str = "") -> None:
    """Create a simple title presentation as a starting point."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0]
        run2.font.size = Pt(22)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    prs.save(output_path)
    print(f"Created: {output_path}")


def read_presentation(input_path: str) -> None:
    """Extract text content from a presentation."""
    from pptx import Presentation

    prs = Presentation(input_path)
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n=== Slide {slide_num} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        print(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    print(" | ".join(cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"[Notes]: {notes}")


def get_info(input_path: str) -> None:
    """Get presentation metadata as JSON."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(input_path)
    info = {
        "slide_count": len(prs.slides),
        "slide_width_inches": round(prs.slide_width / Emu(914400), 2),
        "slide_height_inches": round(prs.slide_height / Emu(914400), 2),
        "slides": [],
    }
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "number": slide_num,
            "shape_count": len(slide.shapes),
            "has_notes": slide.has_notes_slide and bool(
                slide.notes_slide.notes_text_frame.text.strip()
            ),
        }
        info["slides"].append(slide_info)
    print(json.dumps(info, indent=2))


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    cmd = sys.argv[1]
    if cmd == "create":
        create_presentation(
            sys.argv[2],
            sys.argv[3] if len(sys.argv) > 3 else "Untitled Presentation",
            sys.argv[4] if len(sys.argv) > 4 else "",
        )
    elif cmd == "read":
        read_presentation(sys.argv[2])
    elif cmd == "info":
        get_info(sys.argv[2])
    else:
        print(f"Unknown command: {cmd}")
        sys.exit(1)
